from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Brand Colors
ROYAL_BLUE = RGBColor(0x1A, 0x3A, 0x6B)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xFA)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_BLUE = RGBColor(0x2B, 0x5C, 0xA8)

def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height, font_size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_multiline(slide, lines, left, top, width, height, font_size, color=DARK_GRAY, bold=False, spacing=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if spacing:
            p.space_before = Pt(spacing)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"

# ─────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, ROYAL_BLUE)
add_rect(slide, 0, 5.8, 13.33, 0.12, GOLD)

add_text(slide, "LEBC HUTTO", 1, 0.9, 11, 1.2, 52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Digital Media & Growth Strategy", 1, 2.0, 11, 0.8, 26, bold=False, color=GOLD, align=PP_ALIGN.CENTER)
add_text(slide, "Sermon Podcast · Youth Media Network · Paid Advertising · Social Media", 1, 2.8, 11, 0.6, 15, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
add_rect(slide, 4.5, 3.7, 4.33, 0.06, GOLD)
add_text(slide, "Prepared by S2T Designs  |  May 2026", 1, 4.1, 11, 0.5, 13, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
add_text(slide, "Presented to Pastor Arthur L. Spence  |  Little Ebenezer Baptist Church", 1, 4.6, 11, 0.5, 13, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)

add_text(slide, "CONFIDENTIAL — S2T DESIGNS CLIENT PROPOSAL", 1, 6.8, 11, 0.4, 10, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 2 — WHERE YOU ARE TODAY
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, ROYAL_BLUE)
add_text(slide, "WHERE YOU ARE TODAY", 0.4, 0.25, 12, 0.65, 28, bold=True, color=WHITE)
add_rect(slide, 0.4, 0.88, 2.5, 0.06, GOLD)

# 5 stat boxes
stats = [
    ("670", "YouTube Videos\nUploaded Since 2021"),
    ("937", "Facebook\nFollowers"),
    ("17", "Instagram\nFollowers"),
    ("~100", "YouTube\nSubscribers"),
    ("119", "Years of\nFaith & Service"),
]
colors = [ROYAL_BLUE, MEDIUM_BLUE, RGBColor(0xC0, 0x39, 0x2B), ROYAL_BLUE, GOLD]
text_colors = [WHITE, WHITE, WHITE, WHITE, ROYAL_BLUE]

for i, (num, label) in enumerate(stats):
    x = 0.35 + i * 2.55
    add_rect(slide, x, 1.35, 2.25, 1.6, colors[i])
    add_text(slide, num, x, 1.45, 2.25, 0.9, 36, bold=True, color=text_colors[i], align=PP_ALIGN.CENTER)
    add_text(slide, label, x, 2.25, 2.25, 0.7, 11, color=text_colors[i], align=PP_ALIGN.CENTER)

# Strengths vs Gaps
add_rect(slide, 0.35, 3.2, 5.9, 3.9, WHITE)
add_rect(slide, 6.6, 3.2, 6.38, 3.9, WHITE)

add_text(slide, "✅  STRENGTHS", 0.55, 3.3, 5.5, 0.45, 13, bold=True, color=ROYAL_BLUE)
strengths = [
    "✓  Active YouTube — streaming Sunday services since 2021",
    "✓  937 Facebook followers with 100% positive reviews",
    "✓  Active posting — Memorial Day video posted today",
    "✓  Strong 119-year brand story",
    "✓  Multicultural identity — broad community appeal",
    "✓  Youth ministry spanning toddlers to age 25",
]
add_multiline(slide, strengths, 0.55, 3.75, 5.5, 3.0, 11, color=DARK_GRAY, spacing=4)

add_text(slide, "⚠️  GAPS TO CLOSE", 6.8, 3.3, 5.8, 0.45, 13, bold=True, color=RGBColor(0xC0, 0x39, 0x2B))
gaps = [
    "✗  Instagram nearly dormant — only 17 followers",
    "✗  No podcast presence on Spotify or Apple",
    "✗  YouTube not connected to Facebook or Instagram",
    "✗  No Google Business listing — invisible on Maps",
    "✗  No paid advertising — zero new visitor outreach",
    "✗  No structured content calendar or strategy",
]
add_multiline(slide, gaps, 6.8, 3.75, 5.8, 3.0, 11, color=DARK_GRAY, spacing=4)

# ─────────────────────────────────────────────
# SLIDE 3 — THE OPPORTUNITY
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, ROYAL_BLUE)
add_rect(slide, 0, 0, 0.18, 7.5, GOLD)

add_text(slide, "THE OPPORTUNITY", 0.55, 0.3, 12, 0.7, 30, bold=True, color=WHITE)
add_rect(slide, 0.55, 0.95, 3.5, 0.06, GOLD)

add_text(slide, "Hutto is one of the fastest-growing suburbs in Texas.\nLEBC has a 119-year story — and almost no one outside the congregation knows it.", 
         0.55, 1.15, 12.3, 0.9, 15, color=RGBColor(0xCC, 0xDD, 0xFF))

opps = [
    ("📍", "Local Search Visibility", "Families moving to Hutto search 'Baptist church near me' every week. Without a Google Business listing, LEBC is invisible to them."),
    ("🎙️", "Faith Podcast Market", "Christian podcasts are the #1 fastest-growing podcast category. LEBC already has 670 recorded sermons — none of it on Spotify or Apple."),
    ("👶", "Youth Faith Media Gap", "There is virtually no local youth faith content in Williamson County. LEBC's youth ministry (toddlers–25) can fill that gap and attract young families."),
    ("📱", "Social Media Reach", "Facebook's algorithm rewards video and Live content — exactly what LEBC is already producing but not amplifying."),
]

for i, (icon, title, body) in enumerate(opps):
    y = 2.1 + i * 1.2
    add_rect(slide, 0.55, y, 0.65, 0.9, GOLD)
    add_text(slide, icon, 0.55, y + 0.1, 0.65, 0.7, 20, align=PP_ALIGN.CENTER)
    add_text(slide, title, 1.35, y + 0.05, 3.2, 0.4, 13, bold=True, color=GOLD)
    add_text(slide, body, 1.35, y + 0.42, 11.2, 0.55, 11.5, color=RGBColor(0xCC, 0xDD, 0xFF))

# ─────────────────────────────────────────────
# SLIDE 4 — PAID AD CAMPAIGN
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, ROYAL_BLUE)
add_text(slide, "📢  PAID ADVERTISING CAMPAIGN", 0.4, 0.25, 12, 0.65, 26, bold=True, color=WHITE)
add_rect(slide, 0.4, 0.88, 3.5, 0.06, GOLD)

add_text(slide, "Facebook + Instagram | \"Your Church Home in Hutto\"", 0.4, 1.2, 12, 0.45, 14, color=ROYAL_BLUE, bold=True)

# Left col — campaign details
add_rect(slide, 0.35, 1.75, 6.0, 5.3, WHITE)
add_text(slide, "CAMPAIGN DETAILS", 0.6, 1.9, 5.5, 0.4, 12, bold=True, color=ROYAL_BLUE)
details = [
    "🎯  Target Audience: Ages 25–55, 15-mile radius of Hutto TX",
    "      Interests: Church, faith, family, Baptist, community",
    "",
    "💰  Budget: $150–$200/month ($5–7/day)",
    "",
    "📋  Ad Types: Image ads + video carousel",
    "",
    "🎯  Objective: Awareness + page follows + website visits",
    "",
    "📅  Launch: Within 48 hours of approval",
    "",
    "📈  Month 2+: Retargeting visitors with Pastor Spence welcome video",
]
add_multiline(slide, details, 0.6, 2.35, 5.5, 4.4, 11.5, color=DARK_GRAY, spacing=2)

# Right col — ad copy
add_rect(slide, 6.65, 1.75, 6.28, 5.3, ROYAL_BLUE)
add_text(slide, "SAMPLE AD COPY", 6.9, 1.9, 5.8, 0.4, 12, bold=True, color=GOLD)

add_text(slide, "Version A — Awareness", 6.9, 2.35, 5.8, 0.35, 11, bold=True, color=GOLD)
add_text(slide, '"119 years of faith. One community.\nAll are welcome.\n\nLittle Ebenezer Baptist Church has been\nserving Hutto since 1901.\n\nJoin us this Sunday at 11AM\n215 S. Brushy St, Hutto TX\n\n🙏 Little church. Big heart."',
         6.9, 2.72, 5.8, 2.0, 11, color=WHITE, italic=True)

add_rect(slide, 6.9, 4.75, 5.5, 0.04, GOLD)

add_text(slide, "Version B — Visitor Invite", 6.9, 4.85, 5.8, 0.35, 11, bold=True, color=GOLD)
add_text(slide, '"Looking for a church home in Hutto?\nCome experience real worship, community,\nand fellowship at Little Ebenezer.\n\nSunday Service | 11:00 AM\nWe\'d love to meet you. 💙"',
         6.9, 5.22, 5.8, 1.5, 11, color=WHITE, italic=True)

# ─────────────────────────────────────────────
# SLIDE 5 — SERMON PODCAST
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, ROYAL_BLUE)
add_text(slide, "🎙️  THE WORD FROM LITTLE EBENEZER — SERMON PODCAST", 0.4, 0.25, 12.5, 0.65, 22, bold=True, color=WHITE)
add_rect(slide, 0.4, 0.88, 4.0, 0.06, GOLD)

# Key point banner
add_rect(slide, 0.35, 1.2, 12.6, 0.65, GOLD)
add_text(slide, "⚡  ZERO extra work for the church — LEBC already records everything. We just publish it everywhere.", 
         0.55, 1.28, 12.2, 0.5, 13, bold=True, color=ROYAL_BLUE)

# Left — how it works
add_rect(slide, 0.35, 2.0, 6.0, 5.0, WHITE)
add_text(slide, "HOW IT WORKS", 0.6, 2.15, 5.5, 0.4, 12, bold=True, color=ROYAL_BLUE)
steps = [
    "1️⃣   LEBC uploads Sunday service to YouTube",
    "       (You already do this every week)",
    "",
    "2️⃣   S2T downloads audio, cleans it up,",
    "       adds branded intro + outro (30 sec each)",
    "",
    "3️⃣   Episode published to ALL platforms by 3PM Sunday",
    "",
    "4️⃣   60-second sermon clip cut for Instagram Reels",
    "       + Facebook post — drives traffic back to full episode",
    "",
    "5️⃣   Episode is evergreen — searchable forever",
    "       on Spotify and Apple Podcasts",
]
add_multiline(slide, steps, 0.6, 2.6, 5.5, 4.1, 11.5, color=DARK_GRAY, spacing=3)

# Right — distribution
add_rect(slide, 6.65, 2.0, 6.28, 2.4, ROYAL_BLUE)
add_text(slide, "DISTRIBUTED TO", 6.9, 2.15, 5.8, 0.4, 12, bold=True, color=GOLD)
platforms = ["🎵  Spotify", "🍎  Apple Podcasts", "📺  YouTube Music", "🔊  Amazon Music / Alexa", "🎙️  Pocket Casts + iHeart"]
add_multiline(slide, platforms, 6.9, 2.58, 5.8, 1.7, 12, color=WHITE, spacing=3)

add_rect(slide, 6.65, 4.55, 6.28, 2.45, WHITE)
add_text(slide, "FORMAT", 6.9, 4.7, 5.8, 0.4, 12, bold=True, color=ROYAL_BLUE)
fmt = [
    "📌  Host: Pastor Arthur L. Spence",
    "⏱️   Length: 30–50 min (full sermon)",
    "📅  Drops: Every Sunday by 3PM",
    "🎵  Branded intro + outro audio",
    "🔗  Episode notes with scripture references",
]
add_multiline(slide, fmt, 6.9, 5.12, 5.8, 1.7, 11.5, color=DARK_GRAY, spacing=3)

# ─────────────────────────────────────────────
# SLIDE 6 — YOUTH MEDIA NETWORK
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, ROYAL_BLUE)
add_rect(slide, 0, 0, 0.18, 7.5, GOLD)

add_text(slide, "🎙️  LEBC YOUTH MEDIA NETWORK", 0.55, 0.25, 12, 0.65, 28, bold=True, color=WHITE)
add_text(slide, "Three age-specific shows. One powerful brand.", 0.55, 0.88, 12, 0.45, 15, color=GOLD, bold=True)
add_rect(slide, 0.55, 1.3, 12.2, 0.05, GOLD)

shows = [
    {
        "icon": "🧒",
        "name": "\"Little Lights\"",
        "age": "Ages 4–12  |  Fridays",
        "color": RGBColor(0x1A, 0x6B, 0x3A),
        "host": "Youth leader + kid co-hosts",
        "tone": "Playful, storytelling-driven",
        "topics": ["Bible stories retold for kids", "Faith questions answered by kids", "Family devotional content", "Car-ride friendly (8–12 min)"],
        "why": "Parents actively search for faith-based kids content — massive underserved niche on Spotify Kids"
    },
    {
        "icon": "🎓",
        "name": "\"The Upper Room\"",
        "age": "Ages 13–18  |  Wednesdays",
        "color": RGBColor(0x8B, 0x1A, 0x1A),
        "host": "2–3 LEBC teens (rotating)",
        "tone": "Real, honest, peer-to-peer",
        "topics": ["Faith vs. school culture", "Identity + peer pressure", "College decisions + calling", "Dating, boundaries, mental health"],
        "why": "Teens share content made BY other teens. One episode can spread through an entire high school."
    },
    {
        "icon": "🏛️",
        "name": "\"Built Different\"",
        "age": "Ages 19–25  |  Thursdays",
        "color": RGBColor(0x4A, 0x1A, 0x6B),
        "host": "1–2 LEBC young adults",
        "tone": "Mature, real, aspirational",
        "topics": ["Career + calling", "Relationships + faith", "Community service in Hutto", "Guest interviews: local leaders"],
        "why": "Young adults (19–25) are most likely to leave the church — this keeps them connected digitally."
    },
]

for i, show in enumerate(shows):
    x = 0.35 + i * 4.32
    add_rect(slide, x, 1.5, 4.1, 5.6, show["color"])
    add_text(slide, show["icon"], x, 1.6, 4.1, 0.6, 24, align=PP_ALIGN.CENTER)
    add_text(slide, show["name"], x, 2.15, 4.1, 0.45, 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, show["age"], x, 2.57, 4.1, 0.35, 11, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    add_rect(slide, x + 0.2, 2.95, 3.7, 0.04, RGBColor(0xFF, 0xFF, 0xFF))
    add_text(slide, f"Host: {show['host']}", x + 0.15, 3.05, 3.8, 0.35, 10.5, color=WHITE)
    add_text(slide, f"Tone: {show['tone']}", x + 0.15, 3.38, 3.8, 0.35, 10.5, color=WHITE)
    topic_text = "\n".join([f"• {t}" for t in show["topics"]])
    add_text(slide, topic_text, x + 0.15, 3.75, 3.8, 1.4, 10, color=RGBColor(0xCC, 0xDD, 0xFF))
    add_rect(slide, x + 0.15, 5.2, 3.8, 0.04, GOLD)
    add_text(slide, f"💡 {show['why']}", x + 0.15, 5.28, 3.8, 0.85, 9.5, color=GOLD, italic=True)

# ─────────────────────────────────────────────
# SLIDE 7 — YOUR TWO OPTIONS
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, ROYAL_BLUE)
add_text(slide, "YOUR TWO OPTIONS", 0.4, 0.25, 12, 0.65, 28, bold=True, color=WHITE)
add_rect(slide, 0.4, 0.88, 2.5, 0.06, GOLD)

# Option A
add_rect(slide, 0.35, 1.25, 6.1, 5.85, ROYAL_BLUE)
add_text(slide, "OPTION A", 0.6, 1.4, 5.6, 0.45, 18, bold=True, color=GOLD)
add_text(slide, "S2T PRODUCES IT", 0.6, 1.82, 5.6, 0.4, 14, bold=True, color=WHITE)
add_text(slide, "We handle the entire media operation.", 0.6, 2.2, 5.6, 0.35, 11, color=RGBColor(0xCC, 0xDD, 0xFF))
add_rect(slide, 0.6, 2.58, 5.6, 0.04, GOLD)

a_items = [
    ("Sermon Podcast (weekly)", "Included"),
    ("Social Media Mgmt — 2 platforms", "$500/mo"),
    ("Youth Podcast — 1 show", "+$300/mo"),
    ("Youth Podcast — all 3 shows", "+$700/mo"),
    ("Paid Ad Management", "+$300/mo"),
]
for i, (item, price) in enumerate(a_items):
    y = 2.72 + i * 0.52
    add_text(slide, f"• {item}", 0.6, y, 4.2, 0.42, 11, color=WHITE)
    add_text(slide, price, 4.8, y, 1.3, 0.42, 11, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

add_rect(slide, 0.6, 5.42, 5.6, 0.55, GOLD)
add_text(slide, "STARTER PACKAGE: $800/mo + $150–200 ad spend", 0.7, 5.47, 5.4, 0.45, 11, bold=True, color=ROYAL_BLUE)
add_text(slide, "FULL SERVICE: $1,900/mo + ad spend", 0.6, 5.98, 5.6, 0.4, 11, color=RGBColor(0xCC, 0xDD, 0xFF))

# Option B
add_rect(slide, 6.88, 1.25, 6.1, 5.85, WHITE)
add_text(slide, "OPTION B", 7.13, 1.4, 5.6, 0.45, 18, bold=True, color=ROYAL_BLUE)
add_text(slide, "S2T BUILDS IT — LEBC RUNS IT", 7.13, 1.82, 5.6, 0.4, 13, bold=True, color=ROYAL_BLUE)
add_text(slide, "We set everything up and train your team.", 7.13, 2.2, 5.6, 0.35, 11, color=DARK_GRAY)
add_rect(slide, 7.13, 2.58, 5.6, 0.04, ROYAL_BLUE)

b_items = [
    ("Podcast Setup (Spotify + Apple + RSS)", "$400"),
    ("Youth Media Brand Kit (3 shows)", "$350"),
    ("Production SOP + Volunteer Guide", "$150"),
    ("Facebook Ad Campaign Setup", "$250"),
    ("90-min Training Session (Zoom)", "$200"),
]
for i, (item, price) in enumerate(b_items):
    y = 2.72 + i * 0.52
    add_text(slide, f"• {item}", 7.13, y, 4.2, 0.42, 11, color=DARK_GRAY)
    add_text(slide, price, 11.3, y, 1.3, 0.42, 11, bold=True, color=ROYAL_BLUE, align=PP_ALIGN.RIGHT)

add_rect(slide, 7.13, 5.42, 5.6, 0.55, ROYAL_BLUE)
add_text(slide, "ONE-TIME INVESTMENT: $1,350 total", 7.23, 5.47, 5.4, 0.45, 12, bold=True, color=WHITE)
add_text(slide, "Optional: $150/mo check-in retainer", 7.13, 5.98, 5.6, 0.4, 11, color=DARK_GRAY)

# Recommendation banner
add_rect(slide, 0.35, 7.1, 12.63, 0.05, GOLD)

# ─────────────────────────────────────────────
# SLIDE 8 — OUR RECOMMENDATION
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, ROYAL_BLUE)
add_rect(slide, 0, 0, 0.18, 7.5, GOLD)

add_text(slide, "OUR RECOMMENDATION", 0.55, 0.28, 12, 0.65, 28, bold=True, color=WHITE)
add_rect(slide, 0.55, 0.9, 10, 0.06, GOLD)

add_text(slide, "Start with Option B. Upgrade to Option A when you're ready.", 0.55, 1.08, 12, 0.5, 16, color=GOLD, bold=True)
add_text(slide, "This gives LEBC a low-risk entry point with a clear upgrade path — no long-term commitment required.", 
         0.55, 1.58, 12.3, 0.45, 13, color=RGBColor(0xCC, 0xDD, 0xFF))

steps_rec = [
    ("Week 1", "Pastor Spence selects Option A or B\nFacebook page + YouTube admin access granted to S2T"),
    ("Week 2", "Facebook + Instagram ad campaign goes live\nPodcast channel created on Spotify + Apple Podcasts"),
    ("Week 3", "Sermon Podcast Episode 1 published\n(Pulled from existing YouTube library — no new recording)"),
    ("Week 4", "Youth podcast hosts confirmed from LEBC ministry\nShow names + branding finalized"),
    ("Month 2", "First youth episode recorded and published\nYouTube channel optimized (670 videos organized into playlists)"),
    ("Month 3", "All 3 youth shows active\nAd campaign optimized based on first 60 days of data"),
    ("Month 6", "Full 90-day growth report presented to Pastor Spence\nUpgrade decision: stay on Option B or move to Option A"),
]

for i, (week, desc) in enumerate(steps_rec):
    row = i % 4
    col = i // 4
    x = 0.55 + col * 6.45
    y = 2.25 + row * 1.25
    add_rect(slide, x, y, 1.4, 1.0, GOLD)
    add_text(slide, week, x, y + 0.28, 1.4, 0.45, 11, bold=True, color=ROYAL_BLUE, align=PP_ALIGN.CENTER)
    add_rect(slide, x + 1.45, y, 4.75, 1.0, RGBColor(0x1F, 0x45, 0x80))
    add_text(slide, desc, x + 1.6, y + 0.1, 4.5, 0.85, 10.5, color=WHITE)

# ─────────────────────────────────────────────
# SLIDE 9 — 90-DAY TARGETS
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, LIGHT_GRAY)
add_rect(slide, 0, 0, 13.33, 1.1, ROYAL_BLUE)
add_text(slide, "📈  90-DAY SUCCESS TARGETS", 0.4, 0.25, 12, 0.65, 26, bold=True, color=WHITE)
add_rect(slide, 0.4, 0.88, 3.0, 0.06, GOLD)

metrics = [
    ("Facebook Followers", "937", "1,200+", ROYAL_BLUE),
    ("Facebook Engagement Rate", "~0.2%", "5%+", ROYAL_BLUE),
    ("Instagram Followers", "17", "200+", MEDIUM_BLUE),
    ("YouTube Subscribers", "~100", "250+", RGBColor(0xC0, 0x39, 0x2B)),
    ("Sermon Podcast Subscribers", "0", "75+", RGBColor(0x1A, 0x6B, 0x3A)),
    ("Youth Podcast Downloads/ep", "0", "50+", RGBColor(0x8B, 0x1A, 0x1A)),
    ("Google Business Reviews", "0 (no listing)", "15+", RGBColor(0x4A, 0x1A, 0x6B)),
    ("Ad Reach (monthly)", "0", "5,000–10,000", GOLD),
]

# Header row
add_rect(slide, 0.35, 1.25, 5.5, 0.45, ROYAL_BLUE)
add_rect(slide, 5.9, 1.25, 3.3, 0.45, ROYAL_BLUE)
add_rect(slide, 9.25, 1.25, 3.7, 0.45, ROYAL_BLUE)
add_text(slide, "METRIC", 0.55, 1.3, 5.2, 0.38, 11, bold=True, color=WHITE)
add_text(slide, "TODAY", 6.1, 1.3, 3.0, 0.38, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "90-DAY TARGET", 9.45, 1.3, 3.3, 0.38, 11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (metric, now, target, color) in enumerate(metrics):
    y = 1.8 + i * 0.63
    bg = WHITE if i % 2 == 0 else LIGHT_GRAY
    add_rect(slide, 0.35, y, 5.5, 0.55, bg)
    add_rect(slide, 5.9, y, 3.3, 0.55, bg)
    add_rect(slide, 9.25, y, 3.7, 0.55, bg)
    add_text(slide, metric, 0.55, y + 0.08, 5.2, 0.42, 11.5, color=DARK_GRAY)
    add_text(slide, now, 5.9, y + 0.08, 3.3, 0.42, 11.5, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    add_rect(slide, 9.25, y, 3.7, 0.55, color)
    add_text(slide, target, 9.35, y + 0.08, 3.5, 0.42, 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 10 — NEXT STEPS / CLOSE
# ─────────────────────────────────────────────
slide = blank_slide(prs)
add_rect(slide, 0, 0, 13.33, 7.5, ROYAL_BLUE)
add_rect(slide, 0, 0, 0.18, 7.5, GOLD)
add_rect(slide, 0, 6.1, 13.33, 1.4, RGBColor(0x0F, 0x22, 0x44))

add_text(slide, "NEXT STEPS", 0.55, 0.28, 12, 0.65, 30, bold=True, color=WHITE)
add_rect(slide, 0.55, 0.9, 3.0, 0.06, GOLD)
add_text(slide, "What happens after today's meeting:", 0.55, 1.08, 12, 0.4, 14, color=RGBColor(0xCC, 0xDD, 0xFF))

next_steps = [
    ("1", "Select your option", "Option A (S2T produces) or Option B (DIY with S2T setup)\nNo pressure — both paths get you moving immediately"),
    ("2", "Grant page access", "Facebook page admin + YouTube channel manager access to S2T\nThis allows us to launch the ad and optimize your YouTube"),
    ("3", "Identify youth hosts", "Choose 2–3 teens/young adults from your ministry to host the youth shows\nWe'll train them and handle all the technical setup"),
    ("4", "We go to work", "Facebook ad live within 48 hours\nPodcast channel created + Episode 1 published within 2 weeks"),
]

for i, (num, title, desc) in enumerate(next_steps):
    x = 0.55 + (i % 2) * 6.3
    y = 1.7 + (i // 2) * 2.0
    add_rect(slide, x, y, 0.7, 0.7, GOLD)
    add_text(slide, num, x, y + 0.1, 0.7, 0.52, 22, bold=True, color=ROYAL_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, title, x + 0.85, y, 5.1, 0.42, 14, bold=True, color=WHITE)
    add_text(slide, desc, x + 0.85, y + 0.45, 5.1, 1.2, 11, color=RGBColor(0xBB, 0xCC, 0xEE))

add_text(slide, "S2T Designs  |  David Smith", 0.55, 6.2, 8, 0.45, 13, bold=True, color=WHITE)
add_text(slide, "david@s2tdesigns.com  |  s2tdesigns.com", 0.55, 6.58, 8, 0.4, 12, color=RGBColor(0xAA, 0xBB, 0xDD))
add_text(slide, "Let's build something\nLEBC will be proud of. 🙏", 8.5, 6.08, 4.5, 1.2, 18, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

# Save
output_path = "/app/LEBC_Media_Campaign_Proposal.pptx"
prs.save(output_path)
print(f"✅ Deck saved to {output_path}")
